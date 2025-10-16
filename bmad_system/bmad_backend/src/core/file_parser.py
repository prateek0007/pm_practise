"""
File Parser Module for BMAD System

This module handles the parsing of various file types (TXT, MD, JSON, CSV, DOCX, PDF, 
XLSX, PPTX, PNG, JPG, JPEG, GIF, WEBP, MP4, WEBM) for text extraction. It integrates 
with open-source libraries and includes OCR capabilities for image-based PDFs and image files.
"""

import os
import json
import csv
import mimetypes
from typing import Dict, List, Optional, Any, Union
from io import StringIO
from dataclasses import dataclass
from src.utils.logger import get_logger

logger = get_logger(__name__)

@dataclass
class ParsedFile:
    """Represents a parsed file with extracted content"""
    filename: str
    file_type: str
    content: str
    metadata: Dict[str, Any]
    images: List[str] = None  # Paths to extracted images
    
    def __post_init__(self):
        if self.images is None:
            self.images = []

class FileParser:
    """Handles parsing of various file types"""
    
    def __init__(self):
        self.supported_extensions = {
            '.txt': 'text',
            '.md': 'markdown',
            '.json': 'json',
            '.csv': 'csv',
            '.docx': 'word',
            '.pdf': 'pdf',
            '.xlsx': 'excel',
            '.pptx': 'powerpoint',
            '.png': 'image',
            '.jpg': 'image',
            '.jpeg': 'image',
            '.gif': 'image',
            '.webp': 'image',
            '.mp4': 'video',
            '.webm': 'video'
        }
    
    def parse_file(self, file_path: str) -> Optional[ParsedFile]:
        # Routes files to appropriate parser based on file extension
        """
        Parse a file and extract its content
        
        Args:
            file_path: Path to the file to parse
            
        Returns:
            ParsedFile object with extracted content, or None if parsing failed
        """
        try:
            if not os.path.exists(file_path):
                logger.error(f"File not found: {file_path}")
                return None
            
            filename = os.path.basename(file_path)
            file_ext = os.path.splitext(filename)[1].lower()
            
            if file_ext not in self.supported_extensions:
                logger.warning(f"Unsupported file type: {file_ext}")
                return self._parse_as_text(file_path)
            
            file_type = self.supported_extensions[file_ext]
            
            # Route to appropriate parser
            if file_type == 'text':
                return self._parse_text(file_path)
            elif file_type == 'markdown':
                return self._parse_markdown(file_path)
            elif file_type == 'json':
                return self._parse_json(file_path)
            elif file_type == 'csv':
                return self._parse_csv(file_path)
            elif file_type == 'word':
                return self._parse_word(file_path)
            elif file_type == 'pdf':
                return self._parse_pdf(file_path)
            elif file_type == 'excel':
                return self._parse_excel(file_path)
            elif file_type == 'powerpoint':
                return self._parse_powerpoint(file_path)
            elif file_type == 'image':
                return self._parse_image(file_path)
            elif file_type == 'video':
                return self._parse_video(file_path)
            else:
                return self._parse_as_text(file_path)
                
        except Exception as e:
            logger.error(f"Error parsing file {file_path}: {e}")
            return None
    
    def _parse_text(self, file_path: str) -> ParsedFile:
        """Parse plain text file"""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        
        return ParsedFile(
            filename=os.path.basename(file_path),
            file_type='text',
            content=content,
            metadata={'size': os.path.getsize(file_path)}
        )
    
    def _parse_markdown(self, file_path: str) -> ParsedFile:
        """Parse Markdown file"""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        
        return ParsedFile(
            filename=os.path.basename(file_path),
            file_type='markdown',
            content=content,
            metadata={'size': os.path.getsize(file_path)}
        )
    
    def _parse_json(self, file_path: str) -> ParsedFile:
        """Parse JSON file"""
        with open(file_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        
        # Convert JSON to readable string
        content = json.dumps(data, indent=2)
        
        return ParsedFile(
            filename=os.path.basename(file_path),
            file_type='json',
            content=content,
            metadata={'size': os.path.getsize(file_path), 'keys': list(data.keys()) if isinstance(data, dict) else []}
        )
    
    def _parse_csv(self, file_path: str) -> ParsedFile:
        """Parse CSV file"""
        content_lines = []
        metadata = {}
        
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            # Try to detect delimiter
            sample = f.read(1024)
            f.seek(0)
            sniffer = csv.Sniffer()
            delimiter = sniffer.sniff(sample).delimiter
            
            reader = csv.reader(f, delimiter=delimiter)
            rows = list(reader)
            
            if rows:
                headers = rows[0]
                metadata['headers'] = headers
                metadata['row_count'] = len(rows) - 1
                
                # Convert to readable format
                content_lines.append("Headers: " + ", ".join(headers))
                content_lines.append(f"Total rows: {len(rows) - 1}")
                content_lines.append("\\nFirst 5 rows:")
                
                for i, row in enumerate(rows[1:6]):  # First 5 data rows
                    content_lines.append(f"Row {i+1}: " + ", ".join(row))
        
        return ParsedFile(
            filename=os.path.basename(file_path),
            file_type='csv',
            content="\\n".join(content_lines),
            metadata=metadata
        )
    
    def _parse_word(self, file_path: str) -> ParsedFile:
        """Parse Word document using python-docx"""
        try:
            from docx import Document
            
            doc = Document(file_path)
            content_parts = []
            
            # Extract paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    content_parts.append(paragraph.text)
            
            # Extract tables
            for table in doc.tables:
                content_parts.append("\\n[TABLE]")
                for row in table.rows:
                    row_text = " | ".join([cell.text for cell in row.cells])
                    content_parts.append(row_text)
                content_parts.append("[/TABLE]\\n")
            
            content = "\\n".join(content_parts)
            
            return ParsedFile(
                filename=os.path.basename(file_path),
                file_type='word',
                content=content,
                metadata={
                    'size': os.path.getsize(file_path),
                    'paragraphs': len(doc.paragraphs),
                    'tables': len(doc.tables)
                }
            )
            
        except ImportError:
            logger.error("python-docx not installed. Install with: pip install python-docx")
            return self._parse_as_text(file_path)
        except Exception as e:
            logger.error(f"Error parsing Word document: {e}")
            return self._parse_as_text(file_path)
    
    def _parse_pdf(self, file_path: str) -> ParsedFile:
        """Parse PDF using pymupdf"""
        try:
            import fitz  # PyMuPDF
            
            doc = fitz.open(file_path)
            content_parts = []
            images = []
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                
                # Extract text
                text = page.get_text()
                if text.strip():
                    content_parts.append(f"=== Page {page_num + 1} ===")
                    content_parts.append(text)
                
                # Extract images
                image_list = page.get_images()
                for img_index, img in enumerate(image_list):
                    try:
                        xref = img[0]
                        pix = fitz.Pixmap(doc, xref)
                        if pix.n < 5:  # GRAY or RGB
                            img_path = f"/tmp/pdf_image_{page_num}_{img_index}.png"
                            pix.save(img_path)
                            images.append(img_path)
                        pix = None
                    except Exception as e:
                        logger.warning(f"Error extracting image from PDF: {e}")
            
            doc.close()
            content = "\\n".join(content_parts)
            
            return ParsedFile(
                filename=os.path.basename(file_path),
                file_type='pdf',
                content=content,
                metadata={
                    'size': os.path.getsize(file_path),
                    'pages': len(doc),
                    'images_extracted': len(images)
                },
                images=images
            )
            
        except ImportError:
            logger.error("PyMuPDF not installed. Install with: pip install PyMuPDF")
            return self._parse_as_text(file_path)
        except Exception as e:
            logger.error(f"Error parsing PDF: {e}")
            return self._parse_as_text(file_path)
    
    def _parse_excel(self, file_path: str) -> ParsedFile:
        """Parse Excel file using pandas"""
        try:
            import pandas as pd
            
            # Read all sheets
            excel_file = pd.ExcelFile(file_path)
            content_parts = []
            metadata = {'sheets': []}
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                
                content_parts.append(f"=== Sheet: {sheet_name} ===")
                content_parts.append(f"Shape: {df.shape[0]} rows, {df.shape[1]} columns")
                content_parts.append(f"Columns: {', '.join(df.columns.tolist())}")
                
                # Add first few rows
                content_parts.append("\\nFirst 5 rows:")
                content_parts.append(df.head().to_string())
                content_parts.append("\\n")
                
                metadata['sheets'].append({
                    'name': sheet_name,
                    'rows': df.shape[0],
                    'columns': df.shape[1]
                })
            
            content = "\\n".join(content_parts)
            
            return ParsedFile(
                filename=os.path.basename(file_path),
                file_type='excel',
                content=content,
                metadata=metadata
            )
            
        except ImportError:
            logger.error("pandas not installed. Install with: pip install pandas openpyxl")
            return self._parse_as_text(file_path)
        except Exception as e:
            logger.error(f"Error parsing Excel file: {e}")
            return self._parse_as_text(file_path)
    
    def _parse_powerpoint(self, file_path: str) -> ParsedFile:
        """Parse PowerPoint file using python-pptx"""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            content_parts = []
            
            for i, slide in enumerate(prs.slides):
                content_parts.append(f"=== Slide {i + 1} ===")
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        content_parts.append(shape.text)
                
                content_parts.append("\\n")
            
            content = "\\n".join(content_parts)
            
            return ParsedFile(
                filename=os.path.basename(file_path),
                file_type='powerpoint',
                content=content,
                metadata={
                    'size': os.path.getsize(file_path),
                    'slides': len(prs.slides)
                }
            )
            
        except ImportError:
            logger.error("python-pptx not installed. Install with: pip install python-pptx")
            return self._parse_as_text(file_path)
        except Exception as e:
            logger.error(f"Error parsing PowerPoint file: {e}")
            return self._parse_as_text(file_path)
    
    def _parse_image(self, file_path: str) -> ParsedFile:
        """Parse image file using OCR"""
        try:
            # Try to use OCR to extract text from image
            ocr_text = self._extract_text_from_image(file_path)
            
            # Also get image metadata
            metadata = self._get_image_metadata(file_path)
            
            content = f"Image file: {os.path.basename(file_path)}\\n"
            content += f"Dimensions: {metadata.get('width', 'unknown')} x {metadata.get('height', 'unknown')}\\n"
            
            if ocr_text:
                content += f"\\nExtracted text:\\n{ocr_text}"
            else:
                content += "\\nNo text could be extracted from this image."
            
            return ParsedFile(
                filename=os.path.basename(file_path),
                file_type='image',
                content=content,
                metadata=metadata,
                images=[file_path]
            )
            
        except Exception as e:
            logger.error(f"Error parsing image: {e}")
            return ParsedFile(
                filename=os.path.basename(file_path),
                file_type='image',
                content=f"Image file: {os.path.basename(file_path)}\\nError: Could not process image",
                metadata={'size': os.path.getsize(file_path)},
                images=[file_path]
            )
    
    def _parse_video(self, file_path: str) -> ParsedFile:
        """Parse video file (extract metadata only)"""
        try:
            metadata = self._get_video_metadata(file_path)
            
            content = f"Video file: {os.path.basename(file_path)}\\n"
            content += f"Duration: {metadata.get('duration', 'unknown')}\\n"
            content += f"Format: {metadata.get('format', 'unknown')}\\n"
            content += "\\nNote: Video content analysis not implemented."
            
            return ParsedFile(
                filename=os.path.basename(file_path),
                file_type='video',
                content=content,
                metadata=metadata
            )
            
        except Exception as e:
            logger.error(f"Error parsing video: {e}")
            return ParsedFile(
                filename=os.path.basename(file_path),
                file_type='video',
                content=f"Video file: {os.path.basename(file_path)}\\nError: Could not process video",
                metadata={'size': os.path.getsize(file_path)}
            )
    
    def _parse_as_text(self, file_path: str) -> ParsedFile:
        """Fallback: try to parse as text"""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            
            return ParsedFile(
                filename=os.path.basename(file_path),
                file_type='unknown',
                content=content,
                metadata={'size': os.path.getsize(file_path)}
            )
        except Exception as e:
            logger.error(f"Error parsing as text: {e}")
            return ParsedFile(
                filename=os.path.basename(file_path),
                file_type='unknown',
                content=f"Could not parse file: {str(e)}",
                metadata={'size': os.path.getsize(file_path)}
            )
    
    def _extract_text_from_image(self, image_path: str) -> str:
        """Extract text from image using OCR"""
        try:
            import pytesseract
            from PIL import Image
            
            image = Image.open(image_path)
            text = pytesseract.image_to_string(image)
            return text.strip()
            
        except ImportError:
            logger.warning("OCR libraries not installed. Install with: pip install pytesseract Pillow")
            return ""
        except Exception as e:
            logger.error(f"Error extracting text from image: {e}")
            return ""
    
    def _get_image_metadata(self, image_path: str) -> Dict[str, Any]:
        """Get image metadata"""
        try:
            from PIL import Image
            
            with Image.open(image_path) as img:
                return {
                    'width': img.width,
                    'height': img.height,
                    'format': img.format,
                    'mode': img.mode,
                    'size': os.path.getsize(image_path)
                }
        except ImportError:
            return {'size': os.path.getsize(image_path)}
        except Exception as e:
            logger.error(f"Error getting image metadata: {e}")
            return {'size': os.path.getsize(image_path)}
    
    def _get_video_metadata(self, video_path: str) -> Dict[str, Any]:
        """Get video metadata"""
        # This would require ffmpeg-python or similar
        # For now, return basic info
        return {
            'size': os.path.getsize(video_path),
            'format': os.path.splitext(video_path)[1],
            'duration': 'unknown'
        }
    
    def parse_multiple_files(self, file_paths: List[str]) -> List[ParsedFile]:
        """Parse multiple files and return list of ParsedFile objects"""
        results = []
        for file_path in file_paths:
            parsed = self.parse_file(file_path)
            if parsed:
                results.append(parsed)
        return results

